# -*- coding: utf-8 -*-
"""Présentation de soutenance (16:9) — Maintenance prédictive de la corrosion.
Palette topic : acier (sombre) / rouille (accent) / patine (secondaire) / sable (clair).
"""
import os
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "..", "memoire", "figures")

DARK  = RGBColor(0x16, 0x25, 0x2E)
DARK2 = RGBColor(0x1E, 0x33, 0x3E)
RUST  = RGBColor(0xC9, 0x60, 0x1A)
TEAL  = RGBColor(0x2C, 0x7A, 0x7B)
LIGHT = RGBColor(0xF5, 0xF3, 0xEF)
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
INK   = RGBColor(0x20, 0x2B, 0x31)
MUTED = RGBColor(0x6B, 0x7A, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE   = RGBColor(0xCA, 0xDD, 0xDC)
BORD  = RGBColor(0xE3, 0xDE, 0xD4)

HFONT, BFONT = "Georgia", "Calibri"
SW, SH, TOTAL, CT = 13.333, 7.5, 17, 2.3   # CT = haut du contenu (sous titre 2 lignes)

prs = Presentation(); prs.slide_width = I(SW); prs.slide_height = I(SH)
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, color, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, I(x), I(y), I(w), I(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def bg(slide, color):
    rect(slide, -0.06, -0.06, SW + 0.12, SH + 0.12, color)


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         wrap=True, sp_after=4, ls=1.0):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0); p.line_spacing = ls
        for (txt, size, color, bold, italic, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tb


def R(txt, size, color, bold=False, italic=False, font=BFONT):
    return (txt, size, color, bold, italic, font)


def footer(slide, idx, dark=False):
    fg = ICE if dark else MUTED
    text(slide, 0.5, 7.04, 9.0, 0.35, [[R("Maintenance prédictive de la corrosion  ·  ESTL La Salle  ·  2025–2026", 9, fg)]])
    text(slide, 11.4, 7.04, 1.45, 0.35, [[R(f"{idx:02d} / {TOTAL:02d}", 10, RUST, bold=True)]], align=PP_ALIGN.RIGHT)


def accent_bar(slide):
    rect(slide, 0, 0, 0.17, SH, RUST)


def kicker(slide, txt, color=RUST):
    text(slide, 0.92, 0.52, 11.5, 0.35, [[R(txt.upper(), 12.5, color, bold=True)]])


def title(slide, txt, size=31, color=INK, y=0.86):
    text(slide, 0.9, y, 11.7, 1.3, [[R(txt, size, color, bold=True, font=HFONT)]], ls=1.0)


def img_fit(slide, path, bx, by, bw, bh, align="center", valign="middle"):
    p = os.path.join(FIG, path)
    iw, ih = Image.open(p).size
    ar = iw / ih; bar = bw / bh
    if ar > bar:
        w, h = bw, bw / ar
    else:
        h, w = bh, bh * ar
    x = bx + (bw - w) / 2 if align == "center" else (bx if align == "left" else bx + bw - w)
    yy = by + (bh - h) / 2 if valign == "middle" else (by if valign == "top" else by + bh - h)
    slide.shapes.add_picture(p, I(x), I(yy), I(w), I(h))


def card(slide, x, y, w, h, fill=CARD, line=BORD):
    return rect(slide, x, y, w, h, fill, line=line, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def circle(slide, x, y, d, txt, fill=RUST, color=WHITE, size=15):
    c = rect(slide, x, y, d, d, fill, shape=MSO_SHAPE.OVAL)
    tf = c.text_frame; tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color; r.font.name = HFONT
    return c


def new():
    return prs.slides.add_slide(BLANK)


# ═════ 1. TITRE ═════
s = new(); bg(s, DARK); rect(s, 0, 0, 0.28, SH, RUST); rect(s, 0, 6.7, SW, 0.8, DARK2)
text(s, 0.95, 0.66, 11.7, 0.4, [[R("RÉPUBLIQUE DU CAMEROUN  ·  ESTL LA SALLE — GÉNIE INDUSTRIEL & MAINTENANCE", 11, ICE, bold=True)]])
text(s, 0.95, 1.08, 11.5, 0.35, [[R("Mémoire de Master 2 — Maintenance Industrielle et Productique", 13, TEAL, italic=True)]])
text(s, 0.92, 1.75, 11.9, 1.9, [[R("Maintenance prédictive de la corrosion", 44, WHITE, bold=True, font=HFONT)]], ls=1.0)
text(s, 0.95, 3.5, 11.6, 0.9, [[R("Une transition Industrie 3.0 → 4.0 par sonde connectée et apprentissage automatique", 18.5, ICE, italic=True)]], ls=1.1)
rect(s, 0.97, 4.55, 4.7, 0.035, RUST)
text(s, 0.95, 4.8, 9.5, 1.2, [
    [R("Présenté et soutenu par   ", 13, MUTED), R("BATOUMBI IKOND Ricky Parfait", 15, WHITE, bold=True), R("   (Mat. 0111 II17)", 12, MUTED)],
    [R("Sous la supervision de   ", 13, MUTED), R("Dr. TCHAWE", 15, WHITE, bold=True)],
], sp_after=8)
text(s, 0.95, 6.88, 11.0, 0.4, [[R("Année académique 2025 – 2026", 12, ICE, bold=True)]])

# ═════ 2. PLAN ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "Déroulé de la présentation"); title(s, "Plan")
plan = [("1", "Contexte & problématique", "Corrosion, verrou applicatif, transition I3.0 → 4.0"),
        ("2", "Objectifs & démarche", "Quatre objectifs, chaîne ISO 13381-1"),
        ("3", "Méthodologie", "Sonde ER · acquisition IoT · modèle ML · GMAO"),
        ("4", "Résultats", "Validation, performance, régimes, intégration"),
        ("5", "Discussion & apports", "Limites assumées, contributions"),
        ("6", "Conclusion & perspectives", "Bilan et jumeau numérique")]
y = 2.0
for n, t, d in plan:
    circle(s, 0.95, y, 0.6, n, fill=DARK, size=18)
    text(s, 1.76, y - 0.04, 10.6, 0.7, [[R(t, 19, INK, bold=True, font=HFONT)], [R(d, 13, MUTED)]], sp_after=2)
    y += 0.81
footer(s, 2)

# ═════ 3. CONTEXTE ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "Contexte")
title(s, "La corrosion : un enjeu majeur, un verrou applicatif")
stats = [("≈ 2 500", "milliards $ / an", "coût mondial de la corrosion (NACE IMPACT, 2016)"),
         ("1 070", "km de pipeline", "Tchad–Cameroun, cas d'étude (COTCO)"),
         ("30 in", "diamètre nominal", "acier API 5L, corrosion multi-mécanismes")]
x = 0.95
for big, unit, desc in stats:
    card(s, x, 2.2, 3.74, 1.95); rect(s, x, 2.2, 3.74, 0.12, RUST, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + 0.25, 2.46, 3.3, 0.8, [[R(big, 37, DARK, bold=True, font=HFONT)], [R(unit, 13.5, RUST, bold=True)]], sp_after=2)
    text(s, x + 0.25, 3.5, 3.3, 0.6, [[R(desc, 11.5, MUTED)]], ls=1.05)
    x += 3.93
text(s, 0.95, 4.55, 11.5, 2.1, [
    [R("Le verrou n'est plus l'instrumentation.", 17, INK, bold=True, font=HFONT)],
    [R("Les opérateurs disposent déjà de sondes ER en place (cas COTCO) et tracent leurs interventions (SAP). Ce qui manque, c'est l'", 14, INK),
     R("intelligence applicative", 14, RUST, bold=True), R(" : corréler résistance / température / temps, prédire le taux de corrosion et la durée de vie, et relier ces prédictions à l'action.", 14, INK)],
    [R("→ C'est précisément la transition ", 14.5, INK, bold=True), R("Industrie 3.0 → Industrie 4.0", 14.5, TEAL, bold=True),
     R(" que ce travail propose d'opérer, à coût maîtrisé.", 14.5, INK, bold=True)],
], sp_after=9, ls=1.12)
footer(s, 3)

# ═════ 4. PROBLÉMATIQUE ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "Problématique")
title(s, "Un paradoxe : des données riches, sous-exploitées")
gaps = [("Pas de corrélation", "les courbes ER ne sont pas reliées algorithmiquement aux variables de procédé"),
        ("Pas de RUL", "aucune estimation du temps avant défaillance (durée de vie résiduelle)"),
        ("Boucle non prédictive", "les historiques (SAP) existent, mais la prédiction n'y est pas reliée pour générer l'action")]
y = 2.35
for t, d in gaps:
    rect(s, 0.95, y + 0.05, 0.22, 0.22, RUST, shape=MSO_SHAPE.OVAL)
    text(s, 1.35, y, 11.0, 0.6, [[R(t + " — ", 15, INK, bold=True, font=HFONT), R(d, 13.5, INK)]], ls=1.05)
    y += 0.66
card(s, 0.95, 4.5, 11.45, 2.05, fill=DARK, line=DARK)
rect(s, 0.95, 4.5, 0.14, 2.05, RUST, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 1.35, 4.72, 10.8, 0.4, [[R("QUESTION CENTRALE", 12, RUST, bold=True)]])
text(s, 1.35, 5.12, 10.85, 1.35, [[
    R("Dans quelle mesure un système intégré — sonde ER instrumentée (ESP32/HX711/DS18B20), modèle XGBoost à double sortie CR + RUL en protocole ", 14.5, WHITE),
    R("run-to-failure", 14.5, ICE, italic=True),
    R(", et boucle décision → action — permet-il d'opérer la transition I3.0 → 4.0, transposable à COTCO comme aux PME industrielles africaines ?", 14.5, WHITE),
]], ls=1.13)
footer(s, 4)

# ═════ 5. OBJECTIFS ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "Objectifs de l'étude")
title(s, "Quatre objectifs, une chaîne de maintenance prédictive")
chain = ["Détection", "Pronostic", "Diagnostic", "Décision", "Action"]
cx, cw = 0.95, 2.18
for i, c in enumerate(chain):
    col = DARK if i % 2 == 0 else TEAL
    rect(s, cx, 2.28, cw, 0.5, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, cx, 2.35, cw, 0.36, [[R(c, 13, WHITE, bold=True)]], align=PP_ALIGN.CENTER)
    if i < 4:
        text(s, cx + cw - 0.02, 2.29, 0.3, 0.48, [[R("›", 20, RUST, bold=True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += cw + 0.13
text(s, 0.95, 2.86, 11.5, 0.3, [[R("Cadre normatif : ISO 13381-1 (pronostic & durée de vie résiduelle)", 11.5, MUTED, italic=True)]])
os_items = [("OS1", "Concevoir & valider la sonde ER instrumentée IoT", "Montage 2 fils + HX711 24 bits + ESP32, acquisition continue 30 s", TEAL),
            ("OS2", "Entraîner un modèle XGBoost CR + RUL", "Validation leave-one-run-out, interprétabilité SHAP", RUST),
            ("OS3", "Diagnostiquer les régimes & facteurs de variabilité", "Température dominante, répétabilité des conditions, alertes graduées", TEAL),
            ("OS4", "Structurer la boucle décision → action", "Module GMAO maison (ordres de travail, KPIs), transposable open-source", RUST)]
y = 3.32
for tag, t, d, col in os_items:
    card(s, 0.95, y, 11.45, 0.78)
    rect(s, 0.95, y, 1.2, 0.78, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 0.95, y, 1.2, 0.78, [[R(tag, 17, WHITE, bold=True, font=HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.35, y + 0.11, 9.9, 0.6, [[R(t, 14.5, INK, bold=True, font=HFONT)], [R(d, 12, MUTED)]], sp_after=2)
    y += 0.85
footer(s, 5)

# ═════ 6. MÉTHODOLOGIE ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "Méthodologie")
title(s, "Une chaîne complète, du capteur à l'ordre de travail")
img_fit(s, "fig_ii1_architecture.png", 7.15, 2.0, 5.55, 4.55, valign="middle")
steps = [("Capteur", "Sonde ER (fil de fer) + HX711 + DS18B20, piloté par ESP32"),
         ("Données", "Envoi HTTPS → Supabase (PostgreSQL) toutes les 30 s"),
         ("Modèle", "Pipeline Python : nettoyage, compensation thermique, XGBoost CR + RUL, SHAP"),
         ("Décision", "Diagnostic des régimes + alertes graduées (vert / orange / rouge)"),
         ("Action", "Module GMAO maison : ordres de travail automatiques + KPIs")]
y = 2.35
for i, (t, d) in enumerate(steps, 1):
    circle(s, 0.95, y, 0.48, str(i), fill=DARK, size=15)
    text(s, 1.58, y - 0.06, 5.3, 0.85, [[R(t, 14.5, RUST, bold=True, font=HFONT)], [R(d, 12, INK)]], sp_after=2, ls=1.03)
    y += 0.84
footer(s, 6)

# ═════ 7. OS1 SONDE ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "OS1 — Détection", color=TEAL)
title(s, "Sonde à résistance électrique : un montage maison robuste")
img_fit(s, "fig_montage_reel.png", 7.2, 2.35, 5.45, 4.0, valign="top")
text(s, 0.95, 2.35, 5.95, 4.4, [
    [R("Principe — ", 14, INK, bold=True, font=HFONT), R("la résistance d'un fil croît quand la corrosion réduit sa section : ", 13, INK), R("R = ρL/πr².", 13, RUST, bold=True, italic=True)],
    [R("Montage retenu — ", 14, INK, bold=True, font=HFONT), R("2 fils à injection de courant (≈ 1,8 mA), lecture différentielle HX711 24 bits.", 13, INK)],
    [R("Pont de Wheatstone abandonné — ", 14, RUST, bold=True, font=HFONT), R("mode commun du HX711, signal utile noyé dans le bruit.", 13, INK)],
    [R("Valeurs justifiées — ", 14, INK, bold=True, font=HFONT), R("970 Ω : courant faible (évite échauffement / polarisation) mais V lisible (plage ±40 mV).", 13, INK)],
    [R("Température — ", 14, INK, bold=True, font=HFONT), R("DS18B20 protégé (tube + huile) pour un couplage thermique fiable ; cadence 30 s.", 13, INK)],
], sp_after=10, ls=1.1)
footer(s, 7)

# ═════ 8. OS2 MODÈLE ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "OS2 — Pronostic", color=TEAL)
title(s, "XGBoost : le bon modèle pour « beaucoup de points, peu de runs »")
text(s, 0.95, 2.35, 6.0, 0.4, [[R("Pourquoi pas un réseau de neurones ?", 16, INK, bold=True, font=HFONT)]])
text(s, 0.95, 2.85, 6.05, 2.2, [
    [R("Chaque essai ≈ 1 000–2 500 lectures, mais ", 13.5, INK), R("peu d'essais", 13.5, RUST, bold=True), R(" : beaucoup de points, peu de runs.", 13.5, INK)],
    [R("→ MLP / LSTM / Transformer ", 13.5, INK), R("sur-apprennent", 13.5, RUST, bold=True), R(" sur ce volume.", 13.5, INK)],
    [R("→ XGBoost ", 13.5, INK), R("(arbres boostés, tabulaire)", 13.5, TEAL, bold=True), R(" est optimal à faible volume, régularisé, ", 13.5, INK), R("interprétable par SHAP.", 13.5, INK)],
], sp_after=7, ls=1.12)
card(s, 0.95, 5.0, 6.0, 1.55, fill=DARK, line=DARK)
text(s, 1.25, 5.2, 5.5, 0.4, [[R("VALIDATION  leave-one-run-out (LORO)", 12.5, RUST, bold=True)]])
text(s, 1.25, 5.62, 5.5, 0.9, [[R("Un essai entier est retiré puis prédit par les autres : on mesure la ", 12.5, WHITE), R("généralisation à un essai jamais vu", 12.5, ICE, bold=True), R(".", 12.5, WHITE)]], ls=1.1)
img_fit(s, "fig_iii3_features.png", 7.25, 2.25, 5.35, 4.0, valign="middle")
text(s, 7.25, 6.32, 5.35, 0.35, [[R("Importance des variables (SHAP) : dégradation cumulée + température", 11, MUTED, italic=True)]], align=PP_ALIGN.CENTER)
footer(s, 8)

# ═════ 9. RÉSULTATS OS2 ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "Résultats — OS2")
title(s, "Le modèle généralise — là où les conditions sont couvertes")
img_fit(s, "fig_iii3_r2_runs.png", 7.05, 2.45, 5.65, 4.0, valign="middle")
card(s, 0.95, 2.35, 5.8, 1.5)
text(s, 1.2, 2.5, 5.4, 0.5, [[R("R² = +0,29", 38, DARK, bold=True, font=HFONT), R("  (moy. LORO)", 13, MUTED)]])
text(s, 1.2, 3.45, 5.4, 0.4, [[R("Run #12 : +0,50   ·   Run #16 : +0,07   —   bat les baselines", 12.5, RUST, bold=True)]])
text(s, 0.95, 4.15, 5.85, 2.4, [
    [R("Le résultat marquant n'est pas la valeur, mais sa ", 14, INK), R("structure", 14, RUST, bold=True, font=HFONT), R(" :", 14, INK)],
    [R("positif là où les conditions du run testé sont ", 13.5, INK), R("couvertes et répétées", 13.5, TEAL, bold=True), R(", négatif sinon.", 13.5, INK)],
    [R("Ce n'est pas le volume brut qui compte, mais la ", 13.5, INK), R("répétabilité des conditions", 13.5, INK, bold=True), R(".", 13.5, INK)],
], sp_after=8, ls=1.12)
footer(s, 9)

# ═════ 10. OS3 TEMPÉRATURE ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "Résultats — OS3", color=TEAL)
title(s, "Température dominante & deux contre-exemples instructifs")
img_fit(s, "fig_iii4_contrexemples.png", 7.05, 2.45, 5.65, 4.0, valign="middle")
text(s, 0.95, 2.35, 5.85, 1.5, [
    [R("La température est la variable dominante", 15, INK, bold=True, font=HFONT)],
    [R("(Arrhenius : la vitesse de corrosion croît exponentiellement avec T).", 13, INK)],
], sp_after=5, ls=1.1)
text(s, 0.95, 3.55, 5.85, 0.4, [[R("Deux essais imprédictibles = deux facteurs isolés :", 14, INK, bold=True)]])
ce = [("Run #15", "régulation thermique dérivante (31 → 28 °C)"),
      ("Run #17", "concentration d'acide (HCl évaporé avant immersion)")]
y = 4.1
for t, d in ce:
    rect(s, 0.95, y + 0.04, 0.2, 0.2, RUST, shape=MSO_SHAPE.OVAL)
    text(s, 1.28, y, 5.5, 0.55, [[R(t + " — ", 13.5, INK, bold=True), R(d, 12.5, INK)]], ls=1.05)
    y += 0.6
text(s, 0.95, 5.45, 5.85, 1.0, [[R("→ La fiabilité exige le contrôle ", 13.5, INK, bold=True), R("simultané", 13.5, RUST, bold=True), R(" de tous les facteurs, pas seulement la température.", 13.5, INK, bold=True)]], ls=1.1)
footer(s, 10)

# ═════ 11. PHYSIQUE EMBALLEMENT ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "Compréhension physique")
title(s, "Emballement & rupture : une divergence géométrique")
img_fit(s, "fig_iii4_emballement.png", 7.0, 2.3, 5.7, 4.3, valign="middle")
text(s, 0.95, 2.35, 5.85, 4.4, [
    [R("La section varie en r², donc ", 14, INK), R("R s'emballe en 1/r³", 15, RUST, bold=True, italic=True, font=HFONT), R(".", 14, INK)],
    [R("La ", 13.5, INK), R("rupture mécanique coïncide avec le début de l'emballement", 13.5, INK, bold=True), R(" : « section → 0 » = explosion de R ET perte de tenue.", 13.5, INK)],
    [R("Après rupture, ", 13.5, INK), R("la conduction résiduelle", 13.5, TEAL, bold=True), R(" (électrolyte HCl + dernier filament) maintient une résistance croissante jusqu'au circuit ouvert.", 13.5, INK)],
    [R("Cette portion terminale est ", 13.5, INK), R("électrolytique, non métallique", 13.5, RUST, bold=True), R(" : tronquée et exclue de l'apprentissage.", 13.5, INK)],
], sp_after=11, ls=1.13)
footer(s, 11)

# ═════ 12. OS4 GMAO ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "Résultats — OS4", color=TEAL)
title(s, "Boucle décision → action : un module GMAO maison")
text(s, 0.95, 2.35, 11.4, 1.5, [
    [R("Aucun CMMS open-source (GLPI, OpenMaint, Snipe-IT…) n'offre d'API exploitable en version gratuite.", 14, INK)],
    [R("→ ", 14, RUST, bold=True), R("Module GMAO maison", 14.5, INK, bold=True, font=HFONT), R(" intégré à Streamlit + Supabase : à chaque alerte, un ", 14, INK), R("ordre de travail enrichi", 14, RUST, bold=True), R(" (CR, RUL, régime, section) est généré et tracé — sans appel externe.", 14, INK)],
], sp_after=8, ls=1.12)
kpis = [("MTBF · MTTR", "indicateurs de maintenance calculés sur l'historique des OT"),
        ("Gratuit", "coût de licence nul — accessible aux PME africaines"),
        ("Transposable", "même mapping reportable sur un CMMS open-source")]
x = 0.95
for big, d in kpis:
    card(s, x, 4.2, 3.74, 2.25); rect(s, x, 4.2, 3.74, 0.12, TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + 0.28, 4.5, 3.2, 0.7, [[R(big, 23, DARK, bold=True, font=HFONT)]])
    text(s, x + 0.28, 5.35, 3.2, 1.0, [[R(d, 13, MUTED)]], ls=1.12)
    x += 3.93
footer(s, 12)

# ═════ 13. PERSPECTIVE ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "Perspective")
title(s, "Jumeau numérique : prédire puis confirmer")
img_fit(s, "fig_simulateur_run21_nonparam.png", 6.95, 2.0, 5.75, 4.5, valign="middle")
text(s, 0.95, 2.25, 5.8, 1.6, [
    [R("Un jumeau numérique calibré sur les essais réels génère une ", 14, INK), R("bande prédictive de durée de vie", 14, RUST, bold=True, font=HFONT), R(".", 14, INK)],
    [R("Test « prédire-puis-confirmer » sur Run #21 :", 14.5, INK, bold=True)],
], sp_after=6, ls=1.12)
card(s, 0.95, 3.9, 5.8, 1.45, fill=DARK, line=DARK)
text(s, 1.25, 4.1, 5.3, 1.1, [
    [R("Prédit avant l'essai : ", 13.5, ICE), R("13 – 20,5 h", 16, WHITE, bold=True, font=HFONT)],
    [R("Rupture observée : ", 13.5, ICE), R("13,1 h", 16, RUST, bold=True, font=HFONT), R("  ✓ dans la bande", 13.5, WHITE)],
], sp_after=7, ls=1.15)
text(s, 0.95, 5.65, 5.8, 1.0, [[R("Enseignement : on ne prédit bien qu'une morphologie déjà observée ", 13, INK), R("≥ 2 fois", 13, RUST, bold=True), R(" — ce qui quantifie le besoin de répétition.", 13, INK)]], ls=1.12)
footer(s, 13)

# ═════ 14. DISCUSSION ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "Discussion")
title(s, "Limites assumées d'une preuve de concept")
lim = [("Matériau", "fil de fer ≠ acier API 5L : les valeurs absolues de CR ne sont pas directement transposables."),
       ("Effet acide-shunt", "en HCl concentré, l'électrolyte shunte la mesure — d'où l'emploi d'un fil fin."),
       ("Jeu de données", "peu d'essais : métriques bruitées, en cours de consolidation."),
       ("Couverture thermique", "une seule plage (~30 °C) répétée ; ~32 °C encore représentée par un seul essai.")]
y = 2.3
for t, d in lim:
    card(s, 0.95, y, 11.45, 0.92); rect(s, 0.95, y, 0.14, 0.92, RUST, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 1.35, y, 10.8, 0.92, [[R(t + "  —  ", 15, INK, bold=True, font=HFONT), R(d, 13.5, INK)]], ls=1.05, anchor=MSO_ANCHOR.MIDDLE)
    y += 1.02
text(s, 0.95, 6.5, 11.4, 0.45, [[R("Ces limites ne sont pas intrinsèques à la méthode : elles définissent les conditions du passage à l'échelle.", 13.5, TEAL, bold=True, italic=True)]])
footer(s, 14)

# ═════ 15. APPORTS ═════
s = new(); bg(s, LIGHT); accent_bar(s); kicker(s, "Apports du travail")
title(s, "Ce que cette étude apporte")
apf = [("Apport méthodologique", "La fiabilité du ML se lit dans la STRUCTURE du R² (positif si conditions couvertes et répétées) — rarement explicité dans la littérature."),
       ("Chaîne complète low-cost", "De la mesure à l'ordre de travail, en matériel accessible localement (~50 000 FCFA) et logiciels open-source."),
       ("Double transposabilité", "Saut I3.0 → 4.0 essentiellement logiciel chez COTCO ; déploiement autonome pour les PME africaines."),
       ("Rigueur & traçabilité", "Chaque donnée, seuil et décision adossé à une source (norme, référence, loi physique).")]
y = 2.35
for i, (t, d) in enumerate(apf, 1):
    circle(s, 0.95, y + 0.08, 0.55, str(i), fill=TEAL, size=16)
    text(s, 1.7, y, 10.7, 1.0, [[R(t, 16, INK, bold=True, font=HFONT)], [R(d, 13, INK)]], sp_after=3, ls=1.08)
    y += 1.08
footer(s, 15)

# ═════ 16. CONCLUSION ═════
s = new(); bg(s, DARK); rect(s, 0, 0, 0.28, SH, RUST)
text(s, 0.95, 0.66, 11.0, 0.4, [[R("CONCLUSION", 13, RUST, bold=True)]])
text(s, 0.92, 1.12, 11.6, 1.0, [[R("Une brique fondatrice, validée et reproductible", 31, WHITE, bold=True, font=HFONT)]])
concl = [("OS1", "chaîne d'acquisition ER fonctionnelle, suivi d'essais run-to-failure complets."),
         ("OS2", "XGBoost prédit le CR et surpasse les baselines — sous condition de couverture."),
         ("OS3", "température dominante ; la répétabilité des conditions fait la fiabilité."),
         ("OS4", "boucle décision → action démontrée par un module GMAO maison, à coût nul.")]
y = 2.3
for tag, d in concl:
    circle(s, 0.97, y, 0.6, tag[-1], fill=RUST, size=17)
    text(s, 1.8, y + 0.02, 10.6, 0.6, [[R(tag + " — ", 15, RUST, bold=True, font=HFONT), R(d, 14, ICE)]], ls=1.05)
    y += 0.78
text(s, 0.95, 5.7, 11.4, 1.0, [[R("Objectifs partiellement atteints et en voie de consolidation — ", 14.5, WHITE, bold=True),
     R("un prototype doublement transposable, de la mesure à la décision.", 14.5, ICE, italic=True)]], ls=1.15)
footer(s, 16, dark=True)

# ═════ 17. MERCI ═════
s = new(); bg(s, DARK); rect(s, 0, 0, SW, 0.28, RUST); rect(s, 0, SH - 0.28, SW, 0.28, RUST)
text(s, 1.0, 2.5, 11.3, 1.2, [[R("Merci de votre attention", 44, WHITE, bold=True, font=HFONT)]], align=PP_ALIGN.CENTER)
text(s, 1.0, 3.75, 11.3, 0.7, [[R("Questions & discussion", 22, RUST, bold=True, italic=True)]], align=PP_ALIGN.CENTER)
text(s, 1.0, 5.2, 11.3, 0.8, [
    [R("BATOUMBI IKOND Ricky Parfait", 15, ICE, bold=True)],
    [R("Master 2 — Maintenance Industrielle et Productique  ·  ESTL La Salle  ·  2025–2026", 12, MUTED)],
], align=PP_ALIGN.CENTER, sp_after=4)

OUT = os.path.join(HERE, "Soutenance_Corrosion_BATOUMBI.pptx")
prs.save(OUT)
print("OK ->", OUT, "| slides:", len(prs.slides._sldIdLst))
